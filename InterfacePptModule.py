#!/usr/bin/env python 
# -*- coding:utf-8 -*-

# __author__ = 'Hugo Chen'
# __time__   = '2019/1/2'
# God helps who help themselves!

from pptx import Presentation
# 计算PPT尺寸时使用
# 经尝试,正常宽度(4:3)的PPT尺寸为10*7.5(inches)即25.4*19.05(cm)
# 宽屏(16:9)的PPT尺寸为10*5.625(inches)即25.4*14.29(cm)
from pptx.util import Inches
from pptx.util import Pt  # 设置字体时使用
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE  # MSO_SHAPE # 添加形状时使用
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT # PP_ALIGN # 段落设置排列方式
from pptx.enum.text import MSO_VERTICAL_ANCHOR  # 垂直方向的排列方式
from pptx.dml.color import RGBColor # 字体颜色

# 图表使用
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.chart import XL_TICK_MARK
from pptx.enum.chart import XL_DATA_LABEL_POSITION   # XL_LABEL_POSITION
from pptx.enum.chart import XL_LEGEND_POSITION


# 操作接口手册: https://python-pptx.readthedocs.io/en/latest/

'''
PPT file 操作: 新建,打开,保存,查询
'''
# 新建一个PPT:
def createNewPptFile():
    newPpt = Presentation()
    print("slides size:", len(newPpt.slides))
    return newPpt

# 打开PPT:
def openPptFile(path):
    openPpt = Presentation(path)
    return openPpt

# 保存PPT:
def saveAs(openPpt, path):
    openPpt.save(path)

# 获取当前PPT的总页数:
def getSildesNum(pptFile):
    return len(pptFile.slides)

'''
PPT slide 操作: 新建一页
'''
# 添加一页新的PPT: type是新加页的类型,0-10
# 此处type=0: PPT第一页的类型,ppt标题页;
# type=1: 带title和文本框的PPT页;
# type=5: 带title的PPT页;
# type=6: 完全空白页;
def createNewPage(pptFile, type):
    titleSlideLyt = pptFile.slide_layouts[type]
    newSlide = pptFile.slides.add_slide(titleSlideLyt)
    return newSlide

# 通过PPT的页码获取当前Slide:
def getOneSilde(pptFile, num):
    oneSlide = pptFile.slides[num]
    return oneSlide

# 获取PPT中的shapes个数:
def getPageShapesNum(oneSlide):
    return len(oneSlide.shapes)

# 获取PPT中的shapes: (未测试)
def getPageShapes(oneSlide):
    # 在ppt中所有的元素均被当成一个shape,slide.shapes表示幻灯片类中的模型类;
    # placeholders中为每个模型,如采用slide_layouts[1]中包含两个文本框，
    # 所以print len(slide.shapes.placeholders)话为2
    return oneSlide.shapes.placehoders

'''
向PPT中添加内容: 标题栏, 文本框, 形状, 图片, 表格, 以及进行各种设置
'''
'''
文本框类型处理
'''
# 添加标题行titleBox
def addTitleBox(oneSlide, text, left=0.5, top=0.3, width=9, height=1):
    titleBox = oneSlide.shapes.add_textbox(Inches(left), Inches(top),
                                           Inches(width),Inches(height))  # left，top为相对位置，width，height为文本框大小
    textFrame = titleBox.text_frame
    # 此处注意竖直排列方式是text_frame的属性
    textFrame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
    para = textFrame.paragraphs[0]
    # 此处注意水平排列方式是paragraphs[0]的属性
    para.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
    para.text = text
    para.font.size = Pt(25)
    para.font.bold = True
    return titleBox

# 添加文本框textBox
def addTextBox(oneSlide, text, left=0.5, top=1.5, width=9, height=5):
    textBox = oneSlide.shapes.add_textbox(Inches(left), Inches(top), Inches(width),
                                          Inches(height))  # left，top为相对位置，width，height为文本框大小
    textFrame = textBox.text_frame
    textFrame.text = text  # 文本框中文字
    return textBox

# 向textBox添加新的段落,设置textbox中的字体和颜色大小,加粗,斜体,下划线
def addNewPara2TextBox(textBox, newText, size=15, isBold=False, isItalic=False, isUnderline=False, level=1):
    textFrame = textBox.text_frame
    para = textFrame.add_paragraph()  # 在新文本框中添加段落
    para.text = newText  # 段落文字
    para.font.size = Pt(size)  # 文字大小
    para.font.bold = isBold  # 文字加粗
    para.font.italic = isItalic  # 文字斜体
    para.font.underline = isUnderline  # 文字下划线
    para.level = level  # 新段落的级别
    return para

# 设置textbox中的所有段落的字体和颜色大小,加粗,斜体,下划线
# 此处的设置字体删除,改为统一用setShapeFont进行设置:
def setTextFont(textBox, size=15, isBold=False, alignment=PP_PARAGRAPH_ALIGNMENT.LEFT, isItalic=False, isUnderline=False, level=1):
    if not textBox.has_text_frame:
        return
    textFrame = textBox.text_frame
    paraNum = len(textFrame.paragraphs)
    print("一共有多少段:", paraNum)
    for iNum in range(paraNum):
        para = textFrame.paragraphs[iNum]
        para.alignment = alignment  # 靠左排列
        para.font.color.rgb = RGBColor(0x00, 0x00, 0x00)  # 设置文本颜色,默认黑色
        para.font.size = Pt(size)  # 文字大小
        para.font.bold = isBold  # 文字加粗
        para.font.italic = isItalic  # 文字斜体
        para.font.underline = isUnderline  # 文字下划线
        para.level = level  # 新段落的级别

# 向有TitleBox的Slide中添加标题:(首先需要有TitleBox)
def outputTitle(oneSlide, titleTxt):
    title = oneSlide.shapes.title
    title.text = titleTxt

# 向PPT中某Slide的第N个shape里设置文字:(shape的顺序就是按添加到PPT中的顺序来的)
def outputText(newSlide, shapeNum, txt):
    shape = newSlide.shapes[shapeNum]
    shape.text = txt

'''
图片处理类型，包括图片的大小设置
'''
# 添加图片:默认图片添加位置
def addPicture(oneSlide, picPath, isSize=True, left=2, top=2, width=6, height=4.5):
    if isSize == True:
        pic = oneSlide.shapes.add_picture(picPath, Inches(left), Inches(top), Inches(width),
                                          Inches(height))  # 在指定位置按预设值添加图片
    else:
        pic = oneSlide.shapes.add_picture(picPath, Inches(left), Inches(top), width=Inches(width))  # 在指定位置按预设值添加图片,默认图片尺寸
    return pic
'''
形状处理类型，包括形状的大小和填充内容
'''
# 添加形状:默认添加位置，形状和内容
# 形状说明:https://docs.microsoft.com/zh-cn/office/vba/api/Office.MsoAutoShapeType
def addShape(oneSlide, shapeText='', type=MSO_AUTO_SHAPE_TYPE.SQUARE_TABS, left=1, top=2, width=8, height=4):  # 预设位置及大小
    shape = oneSlide.shapes.add_shape(type, Inches(left), Inches(top),
                                      Inches(width), Inches(height))  # 在指定位置按预设值添加类型为PENTAGON的形状
    shape.text = shapeText
    return shape

# 设置shape中的所有段落的字体和颜色大小,加粗,斜体,下划线
def setShapeFont(shape, size=15, isBold=False, isItalic=False, isUnderline=False, level=1):
    if not shape.has_text_frame:
        return
    textFrame = shape.text_frame
    # 设置文字输入框的上下边距
    textFrame.margin_top = Inches(1)  # 上方空出1的空白
    # 设置文字的排版方式
    textFrame.margin_left = 0
    textFrame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE  # 靠中间排列
    textFrame.word_wrap = True
    paraNum = len(textFrame.paragraphs)
    for iNum in range(paraNum):
        para = textFrame.paragraphs[iNum]
        para.font.name = 'Calibri'  # Microsoft Yahei 设置字体
        para.font.size = Pt(size)  # 文字大小
        para.font.color.rgb = RGBColor(0xFF, 0x7F, 0x50)  # 设置文本颜色
        para.font.bold = isBold  # 文字加粗
        para.font.italic = isItalic  # 文字斜体
        para.font.underline = isUnderline  # 文字下划线
        para.level = level  # 新段落的级别

'''
表格处理类型，包括行列尺寸,颜色,字体等
'''
# 添加表格,tables是个表格类型的嵌套list:
def addTables(oneSlide, tableData, left=1, top=2, width=8, height=5):
    if len(tableData) < 0:
        print('no data in the table!')
        return
    rows = len(tableData)
    cols = len(tableData[0])
    table = oneSlide.shapes.add_table(rows, cols, Inches(left), Inches(top),
                                      Inches(width),Inches(height)).table  # 添加表格，并取表格类
    # 遍历将数据写入表格之中:
    for iRow in range(rows):
        for iCol in range(cols):
            # 指定位置写入文本
            table.cell(iRow, iCol).text = tableData[iRow][iCol]
    return table

# 设置表格所有字体格式:
def setTableAllCellFormat(table,size=15, isBold=False, alignment=PP_PARAGRAPH_ALIGNMENT.LEFT, isItalic=False, isUnderline=False, level=1):
    rowNum = len(table.rows)
    colNum = len(table.columns)
    for row in range(rowNum):
        for col in range(colNum):
            cell = table.cell(row, col)
            # if not cell.has_text_frame: # 无此函数
            #     return
            # text_frame.clear()  # 可以删除所有的内容
            textFrame = cell.text_frame
            # 设置文字的排版方式
            textFrame.margin_left = 0
            textFrame.margin_bottom = 0
            textFrame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE  # 靠中间排列
            textFrame.word_wrap = False
            para = textFrame.paragraphs[0]
            para.alignment = alignment
            para.font.name = 'Calibri'  # Microsoft Yahei 设置字体
            para.font.size = Pt(size)  # 文字大小
            para.font.color.rgb = RGBColor(0xFF, 0x7F, 0x50)  # 设置文本颜色
            para.font.bold = isBold  # 文字加粗
            para.font.italic = isItalic  # 文字斜体
            para.font.underline = isUnderline  # 文字下划线
            para.level = level  # 新段落的级别

# 设置表格行格式:
def setTableRowFormat(table,row,size=15, isBold=False, isItalic=False, isUnderline=False, level=1):
    rowNum = len(table.rows)
    if row < rowNum:
        textFrame = table.rows[row].text_frame
        para = textFrame.paragraphs[0]
        para.font.name = 'Calibri'  # Microsoft Yahei 设置字体
        para.font.size = Pt(size)  # 文字大小
        para.font.color.rgb = RGBColor(0xFF, 0x7F, 0x50)  # 设置文本颜色
        para.font.bold = isBold  # 文字加粗
        para.font.italic = isItalic  # 文字斜体
        para.font.underline = isUnderline  # 文字下划线
        para.level = level  # 新段落的级别

# 设置表格列格式:
def setTableColFormat(table,col,size=15, isBold=False, isItalic=False, isUnderline=False, level=1):
    colNum = len(table.columns)
    if col < colNum:
        textFrame = table.columns[col].text_frame
        para = textFrame.paragraphs[0]
        para.font.name = 'Calibri'  # Microsoft Yahei 设置字体
        para.font.size = Pt(size)  # 文字大小
        para.font.color.rgb = RGBColor(0xFF, 0x7F, 0x50)  # 设置文本颜色
        para.font.bold = isBold  # 文字加粗
        para.font.italic = isItalic  # 文字斜体
        para.font.underline = isUnderline  # 文字下划线
        para.level = level  # 新段落的级别

# 设置table的单列宽度
def setTableOneColumnWidth(table, col, width):
    colNum = len(table.columns)
    if col<colNum:
        table.columns[col].width = Inches(width)

# 设置table所有的列宽:
def setTableColumnWidth(table, width):
    colNum = len(table.columns)
    for col in range(colNum):
        table.columns[col].width = Inches(width)

# 设置table的单行行高:
def setTableOneRowHeight(table, row, height):
    rowNum = len(table.rows)
    if row < rowNum:
        table.rows[row].height = Inches(height)

# 设置table所有的行高:
def setTableRowHeight(table, height):
    rowNum = len(table.rows)
    for row in range(rowNum):
        table.rows[row].height = Inches(height)

# 对table进行merge:
def setTableMerge(table, cellRow, cellCol, cellRow2, cellCol2):
    # 1. merge完成后可以判断cell是不是merge过:
    # is_merge_origin返回True是说是开始点,is_spanned返回True是说是扩展的
    # 2. 对merge过的cell使用split()可以恢复原样
    cell = table.cell(cellRow,cellCol)
    cell2 = table.cell(cellRow2,cellCol2)
    cell.merge(cell2)

'''
图表类型处理
'''
# 获取区域位置: title, leftTop, leftBottom, rightTop, rightBottom, titleMiddle
def getPosInSlide(pos):
    if pos == 'title':
        return Inches(0.5), Inches(0.3), Inches(9), Inches(1)
    elif pos == 'middle':
        return Inches(1), Inches(1), Inches(8), Inches(6)
    elif pos == 'leftTop':
        return Inches(0.5), Inches(0.5), Inches(4), Inches(3)
    elif pos == 'leftBottom':
        return Inches(0.5), Inches(4), Inches(4), Inches(3)
    elif pos == 'rightTop':
        return Inches(5.5), Inches(0.5), Inches(4), Inches(3)
    elif pos == 'rightBottom':
        return Inches(5.5), Inches(4), Inches(4), Inches(3)
    elif pos == 'titleMiddle':
        return Inches(0.5), Inches(1.5), Inches(9), Inches(5)
    else:  #
        return Inches(0.5), Inches(1.5), Inches(9), Inches(5)

# 柱状图示例:
def addChartBarDemo(slide):
    # chart1 左上方图
    x, y, cx, cy = getPosInSlide('leftTop')  # 按英尺标准指定x，y值

    chart_data = ChartData()  # 图表data类

    chart_data.categories = [u'A班级得分率', u'B班级得分率']  # 图表加入两栏
    chart_data.add_series(u'得分率对比', (80.5, 60.5))  # 在两栏分别填入数据

    graphic_frame = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy,
        chart_data)  # add_chart(图表类型，xy表示图表位置，cx cy表示图表宽高，并且插入chart_data中规定好的数据）

    chart = graphic_frame.chart  # 从生成的图表中取出图表类
    chart.chart_style = 21  # 图表整体颜色风格

    chart.has_title = True  # 图表是否含有标题，默认为False
    chart.chart_title.text_frame.clear()  # 清除原标题
    new_paragraph = chart.chart_title.text_frame.add_paragraph()  # 添加一行新标题
    new_paragraph.text = '得分率对比'  # 新标题
    new_paragraph.font.size = Pt(15)  # 新标题字体大小

    category_axis = chart.category_axis  # category_axis 为chart的category控制类
    category_axis.has_major_gridlines = True  # 是否显示纵轴线
    category_axis.tick_labels.font.italic = True  # tick_labels为图表下标签，置为斜体
    category_axis.tick_labels.font.size = Pt(12)  # 下标签字体大小
    category_axis.tick_labels.font.color.rgb = RGBColor(255, 0, 0)  # 标签字体颜色

    value_axis = chart.value_axis  # value_axis 为chart的value控制类
    value_axis.maximum_scale = 100.0  # 纵坐标最大值
    value_axis.minimum_scale = 0.0  # 纵坐标最小值
    value_axis.minor_tick_mark = XL_TICK_MARK.CROSS
    value_axis.has_minor_gridlines = True

    tick_labels = value_axis.tick_labels  # tick_labels 为chart的纵轴标签控制类
    tick_labels.number_format = '0%'  # 标签显示样式
    tick_labels.font.bold = True  # 字体加粗
    tick_labels.font.size = Pt(10)  # 字体大小
    tick_labels.font.color.rgb = RGBColor(0, 255, 0)  # 标签颜色

    plot = chart.plots[0]  # 取图表中第一个plot
    plot.has_data_labels = True  # 是否显示数据标签
    data_labels = plot.data_labels  # 数据标签控制类
    data_labels.font.size = Pt(9)  # 字体大小
    data_labels.font.color.rgb = RGBColor(0, 0, 255)  # 字体颜色
    data_labels.position = XL_DATA_LABEL_POSITION.INSIDE_END  # 字体位置


# 饼状图示例:
def addChartPieDemo(slide):
    # chart 2 左下方图
    x, y, cx, cy = getPosInSlide('leftBottom')  # 按英尺标准指定x，y值
    chart_data = ChartData()
    chart_data.categories = ['A', 'B', 'C', 'D']
    chart_data.add_series(u'A班级选项占比', (80, 10, 9, 10))
    chart = slide.shapes.add_chart(XL_CHART_TYPE.PIE, x, y, cx, cy, chart_data).chart  # PIE为饼状图

    chart.has_legend = True  # 是否含有下方的说明
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.horz_offset = 0  # 说明位移量 [-1, 1] 默认为0

    chart.plots[0].has_data_labels = True  # 饼中是否写入数值
    data_labels = chart.plots[0].data_labels
    data_labels.number_format = '0%'  # 数值显示格式
    data_labels.position = XL_DATA_LABEL_POSITION.INSIDE_END  # 数值布局方式

    chart.has_title = True
    chart.chart_title.text_frame.clear()  # 清除原标题
    new_paragraph = chart.chart_title.text_frame.add_paragraph()  # 添加一行新标题
    new_paragraph.text = 'A班级选项占比'  # 新标题
    new_paragraph.font.size = Pt(10)  # 新标题字体大小

# 折线图示例:
def addChartPieDemo2(slide):
    # chart 3 右下方图
    x, y, cx, cy = getPosInSlide('rightBottom')  # 按英尺标准指定x，y值
    chart_data = ChartData()
    chart_data.categories = ['A', 'B', 'C', 'D']
    chart_data.add_series(u'B班级选项占比', (0.1, 0.2, 0.3, 0.4))
    chart = slide.shapes.add_chart(XL_CHART_TYPE.PIE, x, y, cx, cy, chart_data).chart

    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM

    chart.plots[0].has_data_labels = True
    data_labels = chart.plots[0].data_labels
    data_labels.number_format = '0%'
    data_labels.position = XL_DATA_LABEL_POSITION.INSIDE_END

    chart.has_title = True
    chart.chart_title.text_frame.clear()  # 清除原标题
    new_paragraph = chart.chart_title.text_frame.add_paragraph()  # 添加一行新标题
    new_paragraph.text = 'B班级选项占比'  # 新标题
    new_paragraph.font.size = Pt(10)  # 新标题字体大小

# 饼状图示例:
def addChartPieDemo3(slide):
    # chart 4 右上方图
    x, y, cx, cy = getPosInSlide('rightTop')
    chart_data = ChartData()
    chart_data.categories = ['0', '1-3', '4-6', '7-9']
    chart_data.add_series('', (50, 18, 30, 34))
    chart = slide.shapes.add_chart(XL_CHART_TYPE.PIE, x, y, cx, cy, chart_data).chart

    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.font.size = Pt(10)

    chart.plots[0].has_data_labels = True
    data_labels = chart.plots[0].data_labels
    data_labels.number_format = '0%'
    data_labels.position = XL_DATA_LABEL_POSITION.INSIDE_END

    chart.has_title = True
    chart.chart_title.text_frame.clear()
    new_title = chart.chart_title.text_frame.add_paragraph()
    new_title.text = '得分占比'
    new_title.font.size = Pt(10)

# 输出文件：
# def output(newSlide):
#     # prs.slide_layouts中一共预存有1 - 48种，采用第六种为空白幻灯片
#     # slide_layouts[1]为带标题和正文框的ppt，slide_layouts[6]为空白页ppt
#     titleSlideLyt = newPpt.slide_layouts[0]
#     newSlide = newPpt.slides.add_slide(titleSlideLyt)
'''
import win32com.client
import time
import os
ppSaveAsWMV = 37
def savePptWmv(pptSrc, wmvDst):
    ppt = win32com.client.Dispatch('PowerPoint.Application')
    presentation = ppt.Presentations.Open(pptSrc, WithWindow=False)
    presentation.CreateVideo(wmvDst, -1, 4, 720, 24, 60)
    start_time_stamp = time.time()
    while True:
        time.sleep(4)
        try:
            os.rename(wmvDst, wmvDst)
            print
            'success'
            break
        except Exception, e:
            pass
    end_time_stamp = time.time()
    print
    end_time_stamp - start_time_stamp
    ppt.Quit()
    pass

if __name__ == '__main__':
    cover_ppt_to_wmv('d:\\python\\demo.ppt', 'd:\\python\\demo.wmv')
'''
